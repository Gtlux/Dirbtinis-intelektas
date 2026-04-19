from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Spalvų schema ---
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_CARD = RGBColor(0x16, 0x21, 0x3e)
ACCENT = RGBColor(0x00, 0xd2, 0xff)
ACCENT2 = RGBColor(0x7c, 0x3a, 0xed)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0xb0, 0xb0, 0xb0)
GREEN = RGBColor(0x10, 0xb9, 0x81)
ORANGE = RGBColor(0xf5, 0x9e, 0x0b)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
    return slide


def add_content_slide(title, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # Title bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10)
        p.level = 0
        # sub-bullets
        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                ps = tf2.add_paragraph()
                ps.text = "    " + sb
                ps.font.size = Pt(17)
                ps.font.color.rgb = GRAY
                ps.space_before = Pt(4)
                ps.level = 1
    return slide


def add_code_slide(title, code_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # Title
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    # Code box
    code_shape = slide.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.7))
    code_shape.fill.solid()
    code_shape.fill.fore_color.rgb = RGBColor(0x0d, 0x11, 0x17)
    code_shape.line.color.rgb = RGBColor(0x30, 0x36, 0x3d)
    tf2 = code_shape.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.2)
    p = tf2.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xe6, 0xed, 0xf3)
    p.font.name = "Consolas"
    return slide


def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # Title
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    # Table
    cols = len(headers)
    tbl_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(tbl_rows, cols, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.5 * tbl_rows))
    table = table_shape.table
    # Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = ACCENT
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_CARD
    # Rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    return slide


# ============================================================================
# SKAIDRĖS
# ============================================================================

# --- 1. Titulinis ---
add_title_slide(
    "Dirbtinis Intelektas – Paieškos Algoritmai",
    "Užduotys 1.2 | 1.3 | 1.4"
)

# --- 2. Projekto struktūra ---
add_content_slide("Projekto struktūra", [
    "EightPuzzle.py – 8 plytelių dėlionės uždavinys",
    "GraphMapProblem.py – kelio paieška Rumunijos žemėlapyje",
    "farmer.py – fermerio, lapės, vištos, grūdų uždavinys (1.4)",
    "search.py – visi paieškos algoritmai + Problem / Node klasės",
    "maps.py – žemėlapių duomenys",
])

# --- 3. EightPuzzle aprašymas ---
add_content_slide("1.2 – EightPuzzle: kas tai?", [
    "3×3 lentelė su skaičiais 1–8 ir tuščia vieta (0)",
    "Tikslas: stumiant plyteles pasiekti tvarkingą išdėstymą",
    "Pradinė būsena: (2, 4, 3, 1, 5, 6, 7, 8, 0)",
    "Galinė būsena:  (1, 2, 3, 4, 5, 6, 7, 8, 0)",
    "Kintamieji: state (tuple), action (UP/DOWN/LEFT/RIGHT), depth, path_cost",
])

# --- 4. EightPuzzle kintamieji ---
add_table_slide("EightPuzzle – kintamieji", 
    ["Kintamasis", "Tipas", "Aprašymas"],
    [
        ["state", "tuple(9)", "Lentelės pozicijos, 0 = tuščia vieta"],
        ["action", "string", "UP, DOWN, LEFT, RIGHT"],
        ["path_cost", "int", "Žingsnių skaičius nuo pradžios"],
        ["depth", "int", "Gylis paieškos medyje"],
    ])

# --- 5. EightPuzzle būsenų schema ---
add_code_slide("EightPuzzle – sprendimo kelias (BFS)", 
"""Būsena 0 (PRADINĖ)  state=(2,4,3,1,5,6,7,8,0)  depth=0
|2|4|3|    action="UP"
|1|5|6|  ──────────►  Būsena 1: (2,4,3,1,5,0,7,8,6)   |2|4|3|
|7|8| |                                                 |1|5| |
                                                        |7|8|6|
  action="LEFT"
──────────►  Būsena 2: (2,4,3,1,0,5,7,8,6)  |2|4|3|
                                              |1| |5|
                                              |7|8|6|
  action="UP" ► Būsena 3 ► action="LEFT" ► Būsena 4

Būsena 4: (0,2,3,1,4,5,7,8,6)   | |2|3|
                                  |1|4|5|
  action="DOWN"                   |7|8|6|
──────────►  Būsena 5: (1,2,3,0,4,5,7,8,6)  |1|2|3|
                                              | |4|5|
  ► RIGHT ► RIGHT ► DOWN                     |7|8|6|

Būsena 8 (GALINĖ ✅): (1,2,3,4,5,6,7,8,0)  |1|2|3|
                                              |4|5|6|
                                              |7|8| |""")

# --- 6. GraphMapProblem ---
add_content_slide("1.2 – GraphMapProblem: kas tai?", [
    "Kelio paieška tarp miestų Rumunijos žemėlapyje",
    "Grafas su svoriais (atstumai tarp miestų km)",
    "Pradinė būsena: state = \"Arad\"",
    "Galinė būsena:  state = \"Bucharest\"",
    "Kintamieji: state (miesto pavadinimas), action, path_cost, depth",
])

# --- 7. GraphMapProblem būsenų schema ---
add_code_slide("GraphMapProblem – sprendimo kelias (BFS)",
"""Būsena 0: state="Arad",      depth=0, path_cost=0

  ├──► state="Zerind",     depth=1, cost=75
  ├──► state="Sibiu",      depth=1, cost=140
  └──► state="Timisoara",  depth=1, cost=118

       Iš "Sibiu":
       ├──► state="Fagaras",   depth=2, cost=239
       ├──► state="Rimnicu",   depth=2, cost=220
       └──► state="Oradea",    depth=2, cost=291

            Iš "Fagaras":
            └──► state="Bucharest",  depth=3, cost=450  ✅ TIKSLAS!

Sprendimas: Arad → Sibiu → Fagaras → Bucharest""")

# --- 8. BFS ---
add_content_slide("Paieška į plotį (BFS – Breadth-First Search)", [
    "Naudoja FIFO eilę (deque) – pirmas įėjo, pirmas išėjo",
    "frontier.append(child) – deda į GALĄ",
    "frontier.popleft() – ima iš PRIEKIO",
    "Tikrina visas būsenas viename lygyje prieš einant giliau",
    "Garantuoja trumpiausią kelią pagal žingsnių skaičių",
    "Trūkumas: naudoja daug atminties (saugo visą lygį)",
])

# --- 9. DFS ---
add_content_slide("Paieška į gylį (DFS – Depth-First Search)", [
    "Naudoja LIFO steką (list) – paskutinis įėjo, pirmas išėjo",
    "frontier.append(child) – deda į GALĄ",
    "frontier.pop() – ima iš GALO",
    "Eina kuo giliau, kol randa tikslą arba aklavietę",
    "NEGARANTUOJA trumpiausio kelio!",
    "Privalumas: naudoja mažai atminties",
])

# --- 10. BFS vs DFS ---
add_table_slide("BFS vs DFS – palyginimas",
    ["Savybė", "BFS (plotis)", "DFS (gylis)"],
    [
        ["Duomenų struktūra", "FIFO eilė (deque)", "LIFO stekas (list)"],
        ["Elementas imamas", "popleft() – pirmas", "pop() – paskutinis"],
        ["Trumpiausias kelias?", "✅ Taip", "❌ Ne"],
        ["Atminties naudojimas", "Didelis", "Mažas"],
        ["Pilnumas", "✅ Taip", "✅ Taip (su explored)"],
    ])

# --- 11. Euristikos ---
add_content_slide("Euristikos (Heuristics)", [
    "h(n) – apytikslė kaina nuo dabartinės būsenos iki tikslo",
    "EightPuzzle: h = plytelių skaičius NE savo vietoje",
    "GraphMap: h = tiesioginis atstumas (straight-line distance) iki tikslo",
    "Greedy Best-First: f(n) = h(n) – tik euristika",
    "A*: f(n) = g(n) + h(n) – tikra kaina + euristika (OPTIMALUS!)",
    "Uniform Cost: f(n) = g(n) – tik tikra kaina",
])

# --- 12. Architektūra (1.3) ---
add_content_slide("1.3 – AIMA architektūra: klasės", [
    "Problem – abstrakti bazinė klasė (initial, goal, actions, result)",
    "EightPuzzle(Problem) – 3×3 dėlionės realizacija",
    "GraphProblem(Problem) – kelio grafe realizacija",
    "Node – paieškos medžio mazgas (state, parent, action, cost, depth)",
    "Node.expand(problem) – generuoja vaikų mazgus",
    "Node.solution() – grąžina veiksmų seką per parent grandinę",
])

# --- 13. Debug trasavimas ---
add_code_slide("1.3 – EightPuzzle trasavimas (debug)",
"""1. puzzle = EightPuzzle((2,4,3,1,5,6,7,8,0))
   → Problem.__init__(initial=..., goal=(1,2,3,4,5,6,7,8,0))

2. breadth_first_graph_search(puzzle)
   → node = Node(state=(2,4,3,1,5,6,7,8,0), depth=0, cost=0)
   → frontier = deque([node])
   → explored = set()

3. Iteracija 1: node = frontier.popleft()  ← state=(2,4,3,1,5,6,7,8,0)
   → explored.add(state)
   → node.expand(problem):
     → actions() → blank=8, galimi: ['UP', 'LEFT']
     → result(state, 'UP') → sukeičia [8] ir [5] → (2,4,3,1,5,0,7,8,6)
     → child_node() → Node(state=..., parent=node, depth=1, cost=1)
   → goal_test() → False
   → frontier.append(child)

... (BFS tęsiasi lygis po lygio) ...

N. goal_test(child.state) == True → return child
   → .solution() → ['UP','LEFT','UP','LEFT','DOWN','RIGHT','RIGHT','DOWN']""")

# --- 14. Debug GraphMap ---
add_code_slide("1.3 – GraphMapProblem trasavimas (debug)",
"""1. switch_country_map("Romania")
   → maps.romania_map_start = "Arad"
   → maps.romania_map_goal = "Bucharest"
   → maps.romania_map = UndirectedGraph({Arad: {Zerind:75, Sibiu:140, ...}})

2. GraphProblem("Arad", "Bucharest", romania_map)
   → self.initial = "Arad", self.goal = "Bucharest"

3. breadth_first_graph_search(problem)
   → node = Node(state="Arad", depth=0, cost=0)

4. Iteracija 1: popleft() → "Arad"
   → actions("Arad") → graph.get("Arad").keys() → ["Zerind","Sibiu","Timisoara"]
   → child_node("Zerind") → cost = 0 + 75 = 75
   → child_node("Sibiu")  → cost = 0 + 140 = 140
   → frontier = [Zerind, Sibiu, Timisoara]

5. Iteracija 2: popleft() → "Zerind" → vaikai: [Oradea]
6. Iteracija 3: popleft() → "Sibiu"  → vaikai: [Fagaras, Rimnicu]
...
N. expand("Fagaras") → "Bucharest" → goal_test = True! ✅""")

# --- 15. Farmer Problem (1.4) ---
add_content_slide("1.4 – Savo uždavinys: Fermerio problema", [
    "Fermeris turi pervežti per upę: lapę, vištą ir grūdus",
    "Valtyje telpa: fermeris + 1 daiktas",
    "Lapė + višta be fermerio → lapė suės vištą ❌",
    "Višta + grūdai be fermerio → višta suės grūdus ❌",
    "Pradinė būsena: ('L','L','L','L') – visi kairėje",
    "Galinė būsena: ('R','R','R','R') – visi dešinėje",
    "Klasė FarmerProblem(Problem) – naudoja AIMA architektūrą",
])

# --- 16. Farmer sprendimas ---
add_table_slide("1.4 – Fermerio sprendimas (BFS, 7 žingsniai)",
    ["Žingsnis", "Fermeris", "Lapė", "Višta", "Grūdai", "Veiksmas"],
    [
        ["0", "KAIRĖ", "KAIRĖ", "KAIRĖ", "KAIRĖ", "VEZTI_VISTA"],
        ["1", "DEŠINĖ", "KAIRĖ", "DEŠINĖ", "KAIRĖ", "PLAUKTI_VIENAM"],
        ["2", "KAIRĖ", "KAIRĖ", "DEŠINĖ", "KAIRĖ", "VEZTI_LAPE"],
        ["3", "DEŠINĖ", "DEŠINĖ", "DEŠINĖ", "KAIRĖ", "VEZTI_VISTA"],
        ["4", "KAIRĖ", "DEŠINĖ", "KAIRĖ", "KAIRĖ", "VEZTI_GRUDUS"],
        ["5", "DEŠINĖ", "DEŠINĖ", "KAIRĖ", "DEŠINĖ", "PLAUKTI_VIENAM"],
        ["6", "KAIRĖ", "DEŠINĖ", "KAIRĖ", "DEŠINĖ", "VEZTI_VISTA"],
        ["7", "DEŠINĖ", "DEŠINĖ", "DEŠINĖ", "DEŠINĖ", "TIKSLAS! ✅"],
    ])

# --- 17. Farmer Prompt'ai ---
add_content_slide("1.4 – Naudoti Prompt'ai (Google Antigravity)", [
    "Prompt 1:",
    "\"gerai, padaryk dabar sita dali\"",
    "\"1.4 Sukurti savo paieškos uždavinio sistemą naudojant aima",
    "architektūrą bei Google Antigravity (neredaguoti koda,",
    "parodyti atsiskaitymo metu tik Prompt'us)\"",
    "",
    "Prompt 2:",
    "\"gal padaryk geriau kita, kokia fermerio lapes vistos ir",
    "grudu per upe uzdaviny. Tik koda kiekviena eilute aprasyk",
    "ka daro, ir geriau kad koda galima butu pasileist\"",
])

# --- 18. Pabaiga ---
add_title_slide("Ačiū už dėmesį!", "Klausimai?")


# ============================================================================
# Išsaugome
# ============================================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Paieskos_Algoritmai.pptx")
prs.save(output_path)
print(f"Skaidres issaugotos: {output_path}")
